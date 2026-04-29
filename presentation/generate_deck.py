"""
Generates the presentation deck for the Financial Document RAG project.
Creates a clean, professional .pptx with native PowerPoint shapes for diagrams.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Color palette - clean, professional
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
TEAL = RGBColor(0x2C, 0x7D, 0x91)
ORANGE = RGBColor(0xE8, 0x7A, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)


def add_title(slide, text, color=NAVY):
    """Add a title at the top of slide."""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    return title_box


def add_subtitle(slide, text, top=1.0):
    """Add a subtitle below title."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.33), Inches(0.5))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = TEAL
    p.font.italic = True
    return box


def add_text(slide, text, left, top, width, height, size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Add a text box."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
    return box


def add_bullet_list(slide, bullets, left, top, width, height, size=16, color=DARK_GRAY):
    """Add bulleted list."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_box(slide, text, left, top, width, height, fill=TEAL, font_color=WHITE, size=14, bold=True):
    """Add a colored rectangle with text - useful for diagrams."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.text = text
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = font_color
    return shape


def add_arrow(slide, from_x, from_y, to_x, to_y, color=NAVY):
    """Add an arrow connector."""
    line = slide.shapes.add_connector(2,  # Straight line connector
                                       Inches(from_x), Inches(from_y),
                                       Inches(to_x), Inches(to_y))
    line.line.color.rgb = color
    line.line.width = Pt(2.5)
    # Add arrow head
    line_elem = line.line._get_or_add_ln()
    tail_end = etree.SubElement(line_elem, qn('a:tailEnd'))
    tail_end.set('type', 'triangle')
    tail_end.set('w', 'med')
    tail_end.set('h', 'med')
    return line


def add_footer(slide, slide_num, total=15):
    """Add slide number footer."""
    box = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4))
    tf = box.text_frame
    tf.text = f"{slide_num} / {total}"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


# ==================== BUILD PRESENTATION ====================

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]


# -------------------- SLIDE 1: TITLE --------------------
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY

# Main title
box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.33), Inches(1.2))
tf = box.text_frame
tf.text = "Financial Document RAG"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE

# Subtitle
box = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.33), Inches(0.8))
tf = box.text_frame
tf.text = "Architecture Walkthrough & Production Roadmap"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(24)
p.font.color.rgb = ORANGE

# Tagline
box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.33), Inches(0.6))
tf = box.text_frame
tf.text = "Retrieval-Augmented Generation for grounded, citable answers over financial filings"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(16)
p.font.color.rgb = LIGHT_GRAY
p.font.italic = True


# -------------------- SLIDE 2: THE PROBLEM --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "The Problem")

add_text(slide, "Financial analysts drown in dense documents", 0.5, 1.1, 12, 0.5,
         size=20, bold=True, color=TEAL)

# Three problem boxes
add_box(slide, "Document Volume\n\n100+ page 10-Ks\nMultiple companies\nQuarterly updates",
        0.5, 2.0, 4.0, 1.8, fill=NAVY, size=14)
add_box(slide, "Keyword Search Fails\n\n'sales' ≠ 'revenue'\nMisses synonyms\nMisses paraphrases",
        4.7, 2.0, 4.0, 1.8, fill=TEAL, size=14)
add_box(slide, "LLMs Alone Hallucinate\n\nStale training data\nMakes up numbers\nNo citations",
        8.9, 2.0, 4.0, 1.8, fill=ORANGE, size=14)

add_text(slide,
         "RAG bridges this: keep the LLM's reasoning, ground it in real, retrievable, citable documents.",
         0.5, 4.5, 12.33, 0.8, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_bullet_list(slide, [
    "▪ Retrieve only what's relevant - no context window limits",
    "▪ Every answer cites its source - traceability for finance",
    "▪ Update knowledge by re-indexing - no expensive retraining",
], 1.5, 5.5, 11, 1.5, size=15, color=DARK_GRAY)
add_footer(slide, 2)


# -------------------- SLIDE 3: APPROACH --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Approach: Two Pipelines")
add_subtitle(slide, "Separation of concerns is critical for production performance")

# Offline pipeline
add_box(slide, "OFFLINE PIPELINE\n(Run when documents are added)",
        0.5, 1.8, 6.0, 0.7, fill=NAVY, size=14)

add_box(slide, "Parse PDFs", 0.5, 2.7, 1.4, 0.6, fill=TEAL, size=11)
add_arrow(slide, 1.9, 3.0, 2.4, 3.0)
add_box(slide, "Chunk", 2.4, 2.7, 1.0, 0.6, fill=TEAL, size=11)
add_arrow(slide, 3.4, 3.0, 3.9, 3.0)
add_box(slide, "Embed", 3.9, 2.7, 1.0, 0.6, fill=TEAL, size=11)
add_arrow(slide, 4.9, 3.0, 5.4, 3.0)
add_box(slide, "Store", 5.4, 2.7, 1.0, 0.6, fill=TEAL, size=11)

add_text(slide, "Slow, batch-style work", 0.5, 3.5, 6.0, 0.4,
         size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Online pipeline
add_box(slide, "ONLINE PIPELINE\n(Run on every user query)",
        7.0, 1.8, 5.8, 0.7, fill=ORANGE, size=14)

add_box(slide, "Embed Q", 7.0, 2.7, 1.1, 0.6, fill=TEAL, size=11)
add_arrow(slide, 8.1, 3.0, 8.4, 3.0)
add_box(slide, "Search", 8.4, 2.7, 1.0, 0.6, fill=TEAL, size=11)
add_arrow(slide, 9.4, 3.0, 9.7, 3.0)
add_box(slide, "Rerank", 9.7, 2.7, 1.0, 0.6, fill=TEAL, size=11)
add_arrow(slide, 10.7, 3.0, 11.0, 3.0)
add_box(slide, "Generate", 11.0, 2.7, 1.4, 0.6, fill=TEAL, size=11)

add_text(slide, "Fast, real-time work", 7.0, 3.5, 5.8, 0.4,
         size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Key points
add_text(slide, "Why separate?", 0.5, 4.5, 12, 0.4, size=18, bold=True, color=NAVY)
add_bullet_list(slide, [
    "▪ Don't re-parse documents on every query (slow + expensive)",
    "▪ Index once, query many times - classic compute-storage trade-off",
    "▪ Independent scaling: ingestion can be batch (Cloud Batch), queries can be on-demand (Cloud Run)",
    "▪ Different SLAs: ingestion can take minutes; queries must respond in seconds",
], 1.0, 5.0, 11.5, 2.0, size=14)
add_footer(slide, 3)


# -------------------- SLIDE 4: SYSTEM ARCHITECTURE --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "System Architecture")
add_subtitle(slide, "Five core components, each independently swappable")

# Top row: indexing
add_box(slide, "PDF Files", 0.5, 1.7, 1.6, 0.7, fill=DARK_GRAY, size=12)
add_arrow(slide, 2.1, 2.05, 2.4, 2.05)
add_box(slide, "Parser\n(unstructured)", 2.4, 1.7, 1.6, 0.7, fill=TEAL, size=11)
add_arrow(slide, 4.0, 2.05, 4.3, 2.05)
add_box(slide, "Chunker\n(800 chars)", 4.3, 1.7, 1.6, 0.7, fill=TEAL, size=11)
add_arrow(slide, 5.9, 2.05, 6.2, 2.05)
add_box(slide, "Embedder\n(MiniLM)", 6.2, 1.7, 1.6, 0.7, fill=TEAL, size=11)
add_arrow(slide, 7.8, 2.05, 8.1, 2.05)
add_box(slide, "Vector Store\n(ChromaDB)", 8.1, 1.7, 1.8, 0.7, fill=NAVY, size=11)

# Vector store circle - prominent
add_box(slide, "VECTOR\nDATABASE\n\n(ChromaDB +\nHNSW Index)", 5.4, 3.2, 2.5, 1.5, fill=NAVY, size=12)

# Connect indexing to db
add_arrow(slide, 9.0, 2.4, 7.5, 3.2)

# Bottom row: query
add_box(slide, "User Question", 0.5, 5.4, 1.7, 0.7, fill=ORANGE, size=12)
add_arrow(slide, 2.2, 5.75, 2.5, 5.75)
add_box(slide, "Embed Query", 2.5, 5.4, 1.6, 0.7, fill=TEAL, size=11)
add_arrow(slide, 4.1, 5.75, 5.4, 5.0)  # to vector DB
add_arrow(slide, 7.9, 5.0, 8.2, 5.4)   # from vector DB
add_box(slide, "Reranker\n(Cross-Enc)", 8.2, 5.4, 1.7, 0.7, fill=TEAL, size=11)
add_arrow(slide, 9.9, 5.75, 10.2, 5.75)
add_box(slide, "LLM\n(Gemini)", 10.2, 5.4, 1.4, 0.7, fill=ORANGE, size=11)
add_arrow(slide, 11.6, 5.75, 11.9, 5.75)
add_box(slide, "Answer", 11.9, 5.4, 1.2, 0.7, fill=GREEN, size=12)

# Labels
add_text(slide, "INDEXING (offline)", 0.5, 1.3, 6.0, 0.3,
         size=11, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
add_text(slide, "QUERY (online)", 0.5, 5.0, 6.0, 0.3,
         size=11, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

# Key takeaway
add_text(slide,
         "Dependency injection throughout - swap any component without touching the rest",
         0.5, 6.6, 12.33, 0.4, size=14, color=DARK_GRAY,
         align=PP_ALIGN.CENTER, bold=True)
add_footer(slide, 4)


# -------------------- SLIDE 5: INDEXING DEEP DIVE --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Indexing Pipeline: Key Decisions")
add_subtitle(slide, "Where most of the engineering nuance lives")

# Three columns of decisions
add_box(slide, "PARSER", 0.5, 1.5, 4.0, 0.6, fill=NAVY, size=16)
add_bullet_list(slide, [
    "Primary: unstructured library",
    "  → handles complex PDF layouts",
    "Fallback: PyMuPDF (lighter)",
    "  → resilience for prod",
    "Strategy: 'fast' (no OCR)",
    "Tracks page numbers",
    "  → enables citations",
], 0.6, 2.2, 4.0, 4.5, size=12)

add_box(slide, "CHUNKER", 4.7, 1.5, 4.0, 0.6, fill=NAVY, size=16)
add_bullet_list(slide, [
    "Recursive splitting",
    "  → respects natural boundaries",
    "  → ¶, sentence, word, char",
    "800 chars + 200 overlap",
    "  → tested empirically",
    "Page metadata preserved",
    "  → 'NVDA 10-K, page 15'",
], 4.8, 2.2, 4.0, 4.5, size=12)

add_box(slide, "EMBEDDER", 8.9, 1.5, 4.0, 0.6, fill=NAVY, size=16)
add_bullet_list(slide, [
    "all-MiniLM-L6-v2",
    "  → 80MB, 384-dim",
    "Local inference",
    "  → no API costs",
    "  → no data leaves system",
    "Lazy loading pattern",
    "  → fast startup",
], 9.0, 2.2, 4.0, 4.5, size=12)

# Bottom callout
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.6))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_GRAY
box.line.color.rgb = LIGHT_GRAY
tf = box.text_frame
tf.text = "Bug story: original chunker ignored chunk_size param - found via evaluation, fixed it, real comparisons emerged"
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(12)
p.font.italic = True
p.font.color.rgb = DARK_GRAY
add_footer(slide, 5)


# -------------------- SLIDE 6: TWO-STAGE RETRIEVAL --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Retrieval: The Two-Stage Pattern")
add_subtitle(slide, "Standard pattern in production search systems (Bing, Google use this)")

# Stage 1
add_box(slide, "STAGE 1: VECTOR SEARCH (Bi-Encoder)", 0.5, 1.5, 6.0, 0.6, fill=TEAL, size=14)
add_bullet_list(slide, [
    "Query and docs embedded SEPARATELY",
    "Compare with cosine similarity",
    "Fast: 1146 chunks → 20 in milliseconds",
    "Approximate (HNSW algorithm)",
    "Trade-off: speed > precision",
], 0.6, 2.2, 6.0, 2.5, size=13)

# Stage 2
add_box(slide, "STAGE 2: RERANK (Cross-Encoder)", 6.8, 1.5, 6.0, 0.6, fill=ORANGE, size=14)
add_bullet_list(slide, [
    "Query + doc processed TOGETHER",
    "Captures word-level interactions",
    "Slow: rescoring 20 docs adds ~500ms",
    "Much more precise",
    "Trade-off: precision > speed",
], 6.9, 2.2, 6.0, 2.5, size=13)

# Visual flow
add_text(slide, "1146 chunks", 0.5, 5.0, 1.5, 0.5, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_box(slide, "1146", 0.7, 5.5, 1.0, 0.6, fill=DARK_GRAY, size=12)
add_arrow(slide, 1.7, 5.8, 2.4, 5.8)
add_box(slide, "TOP 20", 2.4, 5.4, 1.4, 0.8, fill=TEAL, size=14)
add_arrow(slide, 3.8, 5.8, 4.5, 5.8)
add_box(slide, "TOP 5", 4.5, 5.4, 1.2, 0.8, fill=ORANGE, size=14)

add_text(slide, "All chunks", 0.7, 6.2, 1.0, 0.3, size=10, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Bi-Encoder", 2.4, 6.2, 1.4, 0.3, size=10, color=TEAL, align=PP_ALIGN.CENTER)
add_text(slide, "Cross-Encoder", 4.4, 6.2, 1.4, 0.3, size=10, color=ORANGE, align=PP_ALIGN.CENTER)

# Why this pattern
add_text(slide, "Why both?", 7.0, 5.0, 6.0, 0.4, size=16, bold=True, color=NAVY)
add_bullet_list(slide, [
    "Bi-encoder alone → fast but imprecise",
    "Cross-encoder alone → can't scale to whole DB",
    "Together: best of both worlds",
], 7.0, 5.5, 6.0, 1.5, size=12)
add_footer(slide, 6)


# -------------------- SLIDE 7: GENERATION --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Generation: Preventing Hallucination")
add_subtitle(slide, "Critical for finance - made-up numbers are a non-starter")

# Four prevention strategies
add_box(slide, "1. PROMPT ENGINEERING", 0.5, 1.5, 6.0, 0.7, fill=NAVY, size=14)
add_bullet_list(slide, [
    "\"Use ONLY the provided context\"",
    "Explicit, non-negotiable constraint",
    "Versioned like code",
], 0.6, 2.3, 6.0, 1.5, size=12)

add_box(slide, "2. ALLOW \"I DON'T KNOW\"", 6.8, 1.5, 6.0, 0.7, fill=NAVY, size=14)
add_bullet_list(slide, [
    "Permission to refuse is critical",
    "Better to say nothing than wrong thing",
    "Demo'd in live questions",
], 6.9, 2.3, 6.0, 1.5, size=12)

add_box(slide, "3. LOW TEMPERATURE (0.1)", 0.5, 4.0, 6.0, 0.7, fill=NAVY, size=14)
add_bullet_list(slide, [
    "Deterministic but not stuck",
    "Factual responses, no creativity",
    "(0.7+ for creative writing)",
], 0.6, 4.8, 6.0, 1.5, size=12)

add_box(slide, "4. CITATION REQUIREMENTS", 6.8, 4.0, 6.0, 0.7, fill=NAVY, size=14)
add_bullet_list(slide, [
    "Source labels in every chunk",
    "[nvidia-10k.pdf, page 15]",
    "LLM learns to cite naturally",
], 6.9, 4.8, 6.0, 1.5, size=12)

# Bottom: model choice
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.7))
box.fill.solid()
box.fill.fore_color.rgb = ORANGE
box.line.color.rgb = ORANGE
tf = box.text_frame
tf.text = "Model: Gemini 2.5 Flash  |  Wrapper class makes it trivial to swap to Claude, GPT-4, or any other LLM"
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = WHITE
add_footer(slide, 7)


# -------------------- SLIDE 8: EVALUATION FRAMEWORK --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Evaluation Framework")
add_subtitle(slide, "An ML system without evaluation is just guessing")

# Four metric quadrants
add_box(slide, "RETRIEVAL METRICS", 0.5, 1.5, 6.0, 0.6, fill=TEAL, size=14)
add_bullet_list(slide, [
    "Precision@K: of retrieved chunks, how many are relevant?",
    "Recall@K: of relevant chunks, how many were retrieved?",
    "Compares retrieved chunk IDs vs ground truth",
], 0.6, 2.2, 6.0, 1.5, size=12)

add_box(slide, "GENERATION METRICS", 6.8, 1.5, 6.0, 0.6, fill=TEAL, size=14)
add_bullet_list(slide, [
    "Faithfulness: is answer grounded in context?",
    "Relevance: does it answer the question?",
    "LLM-as-judge for scalable scoring",
], 6.9, 2.2, 6.0, 1.5, size=12)

add_box(slide, "OPERATIONAL METRICS", 0.5, 3.9, 6.0, 0.6, fill=TEAL, size=14)
add_bullet_list(slide, [
    "End-to-end latency (ms)",
    "Per-stage breakdown (retrieval/rerank/LLM)",
    "Throughput, cost per query",
], 0.6, 4.6, 6.0, 1.5, size=12)

add_box(slide, "EXPERIMENT TRACKING", 6.8, 3.9, 6.0, 0.6, fill=TEAL, size=14)
add_bullet_list(slide, [
    "Every run logged with config + results",
    "JSON + CSV export for analysis",
    "Compare experiments via CLI",
], 6.9, 4.6, 6.0, 1.5, size=12)

# Bottom
add_text(slide,
         "Test set: 10 curated questions across NVIDIA & Apple 10-Ks with verified ground truth",
         0.5, 6.4, 12.33, 0.5, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_footer(slide, 8)


# -------------------- SLIDE 9: EXPERIMENT RESULTS --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Experimental Results: Chunk Size Comparison")
add_subtitle(slide, "Data-driven decisions on chunking strategy")

# Table header
header_y = 1.7
add_box(slide, "Chunk Size", 1.0, header_y, 2.0, 0.6, fill=NAVY, size=13)
add_box(slide, "Total Chunks", 3.0, header_y, 2.0, 0.6, fill=NAVY, size=13)
add_box(slide, "Precision@5", 5.0, header_y, 2.0, 0.6, fill=NAVY, size=13)
add_box(slide, "Recall@5", 7.0, header_y, 2.0, 0.6, fill=NAVY, size=13)
add_box(slide, "Latency", 9.0, header_y, 2.0, 0.6, fill=NAVY, size=13)
add_box(slide, "Verdict", 11.0, header_y, 1.5, 0.6, fill=NAVY, size=13)

# Row 1: 400 chars (winner)
y = 2.3
add_box(slide, "400", 1.0, y, 2.0, 0.55, fill=GREEN, size=13)
add_box(slide, "2,134", 3.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "0.340", 5.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "0.567", 7.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "3,046ms", 9.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "WINNER", 11.0, y, 1.5, 0.55, fill=GREEN, size=12)

# Row 2: 800 chars
y = 2.9
add_box(slide, "800", 1.0, y, 2.0, 0.55, fill=TEAL, size=13)
add_box(slide, "1,146", 3.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "0.340", 5.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "0.567", 7.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "3,416ms", 9.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "Tied", 11.0, y, 1.5, 0.55, fill=TEAL, size=12)

# Row 3: 1200 chars
y = 3.5
add_box(slide, "1200", 1.0, y, 2.0, 0.55, fill=RED, size=13)
add_box(slide, "742", 3.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "0.300", 5.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "0.500", 7.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "5,404ms", 9.0, y, 2.0, 0.55, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=13)
add_box(slide, "Worst", 11.0, y, 1.5, 0.55, fill=RED, size=12)

# Insights
add_text(slide, "Key insights:", 0.5, 4.5, 12, 0.4, size=16, bold=True, color=NAVY)
add_bullet_list(slide, [
    "Smaller chunks (400-800) outperform larger chunks on BOTH precision AND latency",
    "1200-char chunks lose context granularity → precision drops 12%",
    "Latency jumps 77% (3s → 5.4s) with larger chunks - more text per LLM call",
    "The bug: original chunker ignored chunk_size param. Fixed → meaningful comparisons emerged",
], 1.0, 5.0, 11.5, 2.0, size=13)
add_footer(slide, 9)


# -------------------- SLIDE 10: LIVE DEMO --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Live Demo")
add_subtitle(slide, "Three questions to showcase: precision, multi-document, graceful failure")

# Three demo cards
add_box(slide, "DEMO 1: SPECIFIC NUMBER", 0.5, 1.7, 4.0, 0.7, fill=NAVY, size=14)
add_bullet_list(slide, [
    "\"What was NVIDIA's revenue for FY2024?\"",
    "",
    "Tests:",
    "  • Specific fact retrieval",
    "  • Number precision",
    "  • Source citation (page 15)",
], 0.6, 2.5, 4.0, 4.0, size=11)

add_box(slide, "DEMO 2: MULTI-DOC", 4.7, 1.7, 4.0, 0.7, fill=TEAL, size=14)
add_bullet_list(slide, [
    "\"Compare risk factors for both companies.\"",
    "",
    "Tests:",
    "  • Cross-document retrieval",
    "  • Synthesis across sources",
    "  • Multiple citations",
], 4.8, 2.5, 4.0, 4.0, size=11)

add_box(slide, "DEMO 3: GRACEFUL FAIL", 8.9, 1.7, 4.0, 0.7, fill=ORANGE, size=14)
add_bullet_list(slide, [
    "\"When was the CEO born?\"",
    "",
    "Tests:",
    "  • Hallucination prevention",
    "  • \"I don't know\" behavior",
    "  • System honesty",
], 9.0, 2.5, 4.0, 4.0, size=11)

# Bottom callout
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.7))
box.fill.solid()
box.fill.fore_color.rgb = NAVY
box.line.color.rgb = NAVY
tf = box.text_frame
tf.text = "→ Switch to Streamlit app now ←"
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = WHITE
add_footer(slide, 10)


# -------------------- SLIDE 11: PRODUCTION & GCP --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Production & GCP Migration")
add_subtitle(slide, "Architecture maps cleanly - swap implementations, keep interfaces")

# Mapping table
header_y = 1.7
add_box(slide, "LOCAL (today)", 1.0, header_y, 5.5, 0.6, fill=DARK_GRAY, size=14)
add_box(slide, "GCP (production)", 6.8, header_y, 5.5, 0.6, fill=TEAL, size=14)

mappings = [
    ("ChromaDB (local files)", "Vertex AI Vector Search"),
    ("Sentence Transformers", "Vertex AI Embeddings API"),
    ("Streamlit (single user)", "Cloud Run (auto-scaling)"),
    ("Gemini API direct", "Vertex AI Gemini (enterprise)"),
    ("Local Python script", "Cloud Batch / Cloud Functions"),
    ("Local logs", "Cloud Logging + Cloud Trace"),
]
y = 2.4
for local, gcp in mappings:
    add_box(slide, local, 1.0, y, 5.5, 0.45, fill=LIGHT_GRAY, font_color=DARK_GRAY, size=12, bold=False)
    add_box(slide, gcp, 6.8, y, 5.5, 0.45, fill=TEAL, font_color=WHITE, size=12, bold=True)
    y += 0.5

# Bottom: extras
add_text(slide, "Plus production essentials:", 0.5, 5.7, 12, 0.4, size=14, bold=True, color=NAVY)
add_bullet_list(slide, [
    "▪ Memorystore for query/embedding caching   ▪ Cloud Monitoring + alerts",
    "▪ A/B testing framework for safe experiments   ▪ Multi-region for latency",
], 1.0, 6.1, 11.5, 1.0, size=12)
add_footer(slide, 11)


# -------------------- SLIDE 12: LESSONS LEARNED --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Lessons Learned")
add_subtitle(slide, "What I'd tell another engineer building this")

# Three lesson cards
add_box(slide, "1", 0.5, 1.7, 0.8, 0.8, fill=ORANGE, size=24)
add_text(slide, "Build Evaluation FIRST", 1.5, 1.7, 11, 0.5, size=20, bold=True, color=NAVY)
add_text(slide, "Without metrics, you can't tell if changes are improvements. The chunker bug only surfaced because I was running rigorous experiments. Evaluation is not optional - it's how you know your system actually works.",
         1.5, 2.2, 11, 1.2, size=13, color=DARK_GRAY)

add_box(slide, "2", 0.5, 3.6, 0.8, 0.8, fill=ORANGE, size=24)
add_text(slide, "Two-Stage Retrieval is Production-Grade", 1.5, 3.6, 11, 0.5, size=20, bold=True, color=NAVY)
add_text(slide, "The bi-encoder + cross-encoder pattern isn't optional for serious RAG. Skipping reranking is the difference between mediocre and good results. This is how every major search system works.",
         1.5, 4.1, 11, 1.2, size=13, color=DARK_GRAY)

add_box(slide, "3", 0.5, 5.5, 0.8, 0.8, fill=ORANGE, size=24)
add_text(slide, "Prompt Engineering IS Engineering", 1.5, 5.5, 11, 0.5, size=20, bold=True, color=NAVY)
add_text(slide, "Treat the system prompt as critical, versioned infrastructure. Explicit constraints, fallback behaviors, citation requirements - this is what separates a hallucinating system from a trustworthy one.",
         1.5, 6.0, 11, 1.2, size=13, color=DARK_GRAY)
add_footer(slide, 12)


# -------------------- SLIDE 13: WHAT I'D DO NEXT --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "What I'd Do Next")
add_subtitle(slide, "Roadmap if I had another sprint on this")

# Five roadmap items
roadmap = [
    ("Hybrid Retrieval", "Combine semantic search with BM25 keyword search for better coverage on financial terminology and exact term matches"),
    ("Query Rewriting", "Use a smaller LLM to expand or clarify ambiguous queries before retrieval - improves recall significantly"),
    ("Caching Layer", "Common questions are common. Cache (query → embedding) and (query → answer) at multiple levels"),
    ("Domain Fine-Tuning", "Fine-tune the reranker on financial Q&A data for better precision on domain-specific terminology"),
    ("Multi-Modal Support", "Financial docs have charts and tables that pure text extraction misses - integrate vision models"),
]

y = 1.6
for i, (title, desc) in enumerate(roadmap, 1):
    add_box(slide, str(i), 0.5, y, 0.6, 0.6, fill=TEAL, size=18)
    add_text(slide, title, 1.3, y, 4.0, 0.5, size=15, bold=True, color=NAVY)
    add_text(slide, desc, 1.3, y + 0.4, 11.5, 0.6, size=11, color=DARK_GRAY)
    y += 1.05
add_footer(slide, 13)


# -------------------- SLIDE 14: KNOWN LIMITATIONS --------------------
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "Known Limitations")
add_subtitle(slide, "Honesty about what's not yet there")

limitations = [
    ("Single-Tenant", "Not multi-user safe - would need user-scoped collections and access control",
     "Solution: collection-per-user pattern in Vertex AI Vector Search"),
    ("No Table Extraction", "Financial docs have tables (income statements, balance sheets) that pure text extraction handles poorly",
     "Solution: integrate Camelot or unstructured's table extraction"),
    ("No Conversation History", "Each query is independent - can't ask follow-ups like 'and for last quarter?'",
     "Solution: session state with conversation memory in Cloud Firestore"),
    ("Small Eval Set", "10 curated questions; production needs hundreds with continuous growth",
     "Solution: LLM-assisted ground truth generation + human review pipeline"),
    ("No Production Observability", "Local logs only - no metrics, traces, or alerts",
     "Solution: Cloud Monitoring + Cloud Trace + custom dashboards"),
]

y = 1.7
for title, problem, solution in limitations:
    add_box(slide, title, 0.5, y, 2.5, 0.6, fill=RED, size=12)
    add_text(slide, problem, 3.2, y, 6.5, 0.3, size=11, bold=True, color=DARK_GRAY)
    add_text(slide, "→ " + solution, 3.2, y + 0.3, 9.5, 0.3, size=10, color=TEAL)
    y += 1.0
add_footer(slide, 14)


# -------------------- SLIDE 15: CLOSING --------------------
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY

# Centered title
box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.33), Inches(1.0))
tf = box.text_frame
tf.text = "Thank You"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE

# Summary
box = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.33), Inches(1.0))
tf = box.text_frame
tf.text = "End-to-end RAG system with rigorous evaluation and production roadmap"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(20)
p.font.color.rgb = ORANGE

# Q&A invitation
box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.0))
tf = box.text_frame
tf.text = "Ready for code & design questions"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(28)
p.font.color.rgb = WHITE

box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(11.33), Inches(0.6))
tf = box.text_frame
tf.text = "Architecture · Trade-offs · Implementation · Production · Anything"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = LIGHT_GRAY


# ==================== SAVE ====================
output_path = "presentation/financial_rag_deck.pptx"
prs.save(output_path)
print(f"Deck generated: {output_path}")
print(f"Total slides: {len(prs.slides)}")
